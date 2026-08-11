from pathlib import Path

from app.domain.chunk import Block

# Legacy binary Office formats (.doc, .ppt) need a converter (e.g. LibreOffice
# headless) neither python-docx nor python-pptx can read -- that's real infra
# Step 2 deliberately doesn't add yet. Raising here is correct, not a gap to
# silently paper over: ingest_source() catches it and degrades the *file*
# (P5), it just can't be parsed into blocks at all.
_UNSUPPORTED_LEGACY = {".doc", ".ppt"}

# No OCR in Step 2 -- an image-only upload (e.g. a screenshot of a deck slide)
# always yields zero blocks, which ingest_source() treats as "no extractable
# text" and degrades. Deliberate scope cut, not an oversight.
_IMAGE_TYPES = {".png", ".jpg", ".jpeg"}


def parse_by_type(file_path: str, source_kind: str) -> list[Block]:
    ext = Path(file_path).suffix.lower()
    if ext in _UNSUPPORTED_LEGACY:
        raise ValueError(f"unsupported legacy format {ext} -- convert to {ext}x first")
    if ext in _IMAGE_TYPES:
        return []
    if ext == ".pdf":
        return _blocks_from_paragraphs(_extract_pdf(file_path))
    if ext == ".docx":
        return _blocks_from_paragraphs(_extract_docx(file_path))
    if ext == ".pptx":
        return _blocks_from_pptx(file_path)
    if ext == ".csv":
        return _blocks_from_tabular(file_path, sheet=None)
    if ext == ".xlsx":
        return _blocks_from_tabular(file_path, sheet=0)
    return _blocks_from_paragraphs(Path(file_path).read_text(encoding="utf-8", errors="ignore"))


def extract_text(file_path: str) -> str:
    """Kept for the Step 1 call sites that still want flat text rather than
    Block-level structure."""
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(file_path)
    if ext == ".docx":
        return _extract_docx(file_path)
    return Path(file_path).read_text(encoding="utf-8", errors="ignore")


def _extract_pdf(file_path: str) -> str:
    import pymupdf

    with pymupdf.open(file_path) as doc:
        return "\n\n".join(page.get_text() for page in doc)


def _extract_docx(file_path: str) -> str:
    import docx

    document = docx.Document(file_path)
    return "\n\n".join(p.text for p in document.paragraphs)


def split_paragraphs(text: str) -> list[str]:
    """Structural (L1) split: one chunk per non-empty paragraph, in document order."""
    return [p.strip() for p in text.split("\n\n") if p.strip()]


def _blocks_from_paragraphs(text: str) -> list[Block]:
    paragraphs = split_paragraphs(text)
    return [
        Block(text=p, start=i, end=i, is_structural=True, confidence=1.0)
        for i, p in enumerate(paragraphs)
    ]


def _blocks_from_pptx(file_path: str) -> list[Block]:
    from pptx import Presentation

    prs = Presentation(file_path)
    blocks: list[Block] = []
    for i, slide in enumerate(prs.slides):
        texts = [
            shape.text_frame.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        slide_text = "\n".join(texts)
        if not slide_text:
            continue  # image-only slide -- omitted, not a degraded empty block
        blocks.append(Block(text=slide_text, start=i, end=i, is_structural=True, confidence=1.0))
    return blocks


def _blocks_from_tabular(file_path: str, sheet: int | None) -> list[Block]:
    import pandas as pd

    df = pd.read_csv(file_path) if sheet is None else pd.read_excel(file_path, sheet_name=sheet)
    blocks: list[Block] = []
    for i, row in df.iterrows():
        cells = [f"{col}: {val}" for col, val in row.items() if pd.notna(val)]
        if not cells:
            continue
        blocks.append(Block(text=", ".join(cells), start=i, end=i, is_structural=True, confidence=1.0))
    return blocks
