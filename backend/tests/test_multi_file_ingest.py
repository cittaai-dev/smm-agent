from pathlib import Path

import pandas as pd
from pptx import Presentation

from app.infra.db import get_session
from app.workers.ingest import ingest_file


def _chunk_count(brand_id: str) -> int:
    with get_session() as session:
        row = session.execute(
            "SELECT count(*) AS n FROM chunk WHERE kb_id = :kb", {"kb": f"run:{brand_id}"}
        ).mappings().first()
    return row["n"]


def _source_file_status(brand_id: str) -> list[str]:
    with get_session() as session:
        rows = session.execute(
            "SELECT status FROM source_file WHERE brand_id = :b", {"b": brand_id}
        ).mappings().all()
    return [r["status"] for r in rows]


def _make_pptx(path: Path, slide_texts: list[str | None]) -> str:
    prs = Presentation()
    layout = prs.slide_layouts[5]  # blank-ish layout with a title placeholder
    for text in slide_texts:
        slide = prs.slides.add_slide(layout)
        if text is not None:
            box = slide.shapes.add_textbox(0, 0, 100, 100)
            box.text_frame.text = text
    prs.save(str(path))
    return str(path)


def _make_csv(path: Path) -> str:
    pd.DataFrame(
        {"metric": ["followers", "engagement_rate"], "value": [12000, 0.034]}
    ).to_csv(path, index=False)
    return str(path)


def _make_xlsx(path: Path) -> str:
    pd.DataFrame(
        {"platform": ["Instagram", "TikTok"], "monthly_reach": [50000, 120000]}
    ).to_excel(path, index=False)
    return str(path)


def test_all_file_types_ingest(tmp_path):
    files = {
        "pptx": _make_pptx(tmp_path / "brand_deck.pptx", ["Acme Store overview and mission."]),
        "csv": _make_csv(tmp_path / "insights.csv"),
        "xlsx": _make_xlsx(tmp_path / "analytics.xlsx"),
    }
    for kind, path in files.items():
        result = ingest_file(brand_id=f"brand-{kind}", file_path=path)
        assert result == "ingested", f"{kind} failed: {result}"
        assert _chunk_count(f"brand-{kind}") > 0


def test_scanned_image_only_pptx_degrades_not_fails(tmp_path):
    path = _make_pptx(tmp_path / "scanned_no_text.pptx", [None, None])

    result = ingest_file(brand_id="brand-x", file_path=path)

    assert result == "degraded"
    assert _chunk_count("brand-x") == 0
    assert _source_file_status("brand-x") == ["degraded"]


def test_unsupported_legacy_format_degrades_not_fails(tmp_path):
    path = tmp_path / "old_brand_deck.ppt"
    path.write_bytes(b"not a real legacy ppt file")

    result = ingest_file(brand_id="brand-legacy", file_path=str(path))

    assert result == "degraded"
    assert _chunk_count("brand-legacy") == 0


def test_reupload_across_file_types_is_noop(tmp_path):
    path = _make_csv(tmp_path / "insights.csv")

    ingest_file(brand_id="brand-x", file_path=path)
    n1 = _chunk_count("brand-x")

    ingest_file(brand_id="brand-x", file_path=path)
    n2 = _chunk_count("brand-x")

    assert n1 == n2
    assert n1 > 0


def test_source_kind_inferred_from_extension(tmp_path):
    csv_path = _make_csv(tmp_path / "insights.csv")
    ingest_file(brand_id="brand-x", file_path=csv_path)

    with get_session() as session:
        row = session.execute(
            "SELECT source_kind FROM document_registry WHERE kb_id = :kb", {"kb": "run:brand-x"}
        ).mappings().first()
    assert row["source_kind"] == "analytics"
